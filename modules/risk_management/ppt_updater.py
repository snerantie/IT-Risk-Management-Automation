"""
PowerPoint Updater for Risk Management
Updates PowerPoint slides with risk data and charts
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
from typing import Dict, Optional
import pandas as pd


class RiskPowerPointUpdater:
    """Updates PowerPoint slides with risk management data"""
    
    def __init__(self, template_path: str, output_path: str, logger=None):
        """
        Initialize the PowerPoint updater
        
        Args:
            template_path: Path to PowerPoint template
            output_path: Path to save updated presentation
            logger: Logger instance
        """
        self.template_path = Path(template_path)
        self.output_path = Path(output_path)
        self.logger = logger
        self.presentation = None
    
    def _log(self, message, level="info"):
        """Log message if logger is available"""
        if self.logger:
            getattr(self.logger, level)(message)
        else:
            print(f"{level.upper()}: {message}")
    
    def load_template(self) -> bool:
        """
        Load PowerPoint template
        
        Returns:
            True if successful, False otherwise
        """
        try:
            if not self.template_path.exists():
                self._log(f"Template not found: {self.template_path}", "error")
                return False
            
            self._log(f"Loading PowerPoint template: {self.template_path}")
            self.presentation = Presentation(str(self.template_path))
            self._log(f"Template loaded successfully ({len(self.presentation.slides)} slides)")
            return True
            
        except Exception as e:
            self._log(f"Error loading template: {str(e)}", "error")
            return False
    
    def update_risk_slide(
        self,
        slide_number: int,
        data: pd.DataFrame,
        charts: Dict[str, str],
        statistics: Dict
    ) -> bool:
        """
        Update risk management slide with data and charts
        
        Args:
            slide_number: Slide index (0-based)
            data: Risk data DataFrame
            charts: Dictionary of chart names to file paths
            statistics: Summary statistics
            
        Returns:
            True if successful, False otherwise
        """
        try:
            if self.presentation is None:
                self._log("Presentation not loaded!", "error")
                return False
            
            if slide_number >= len(self.presentation.slides):
                self._log(f"Slide {slide_number} does not exist!", "error")
                return False
            
            self._log(f"Updating slide {slide_number} with risk data...")
            
            slide = self.presentation.slides[slide_number]
            
            # Add charts to slide
            self._add_charts_to_slide(slide, charts)
            
            # Add summary text
            self._add_summary_text(slide, statistics)
            
            self._log("Risk slide updated successfully")
            return True
            
        except Exception as e:
            self._log(f"Error updating risk slide: {str(e)}", "error")
            return False
    
    def _add_charts_to_slide(self, slide, charts: Dict[str, str]):
        """
        Add charts to slide
        
        Args:
            slide: PowerPoint slide object
            charts: Dictionary of chart names to file paths
        """
        # Define positions for charts (adjust as needed)
        chart_positions = {
            'stacked_bar': {'left': Inches(0.5), 'top': Inches(1.5), 'width': Inches(9), 'height': Inches(3)},
            'pie_status': {'left': Inches(0.5), 'top': Inches(5), 'width': Inches(4), 'height': Inches(2.5)},
            'pie_categories': {'left': Inches(5.5), 'top': Inches(5), 'width': Inches(4), 'height': Inches(2.5)}
        }
        
        for chart_name, chart_path in charts.items():
            if chart_name in chart_positions and Path(chart_path).exists():
                pos = chart_positions[chart_name]
                slide.shapes.add_picture(
                    chart_path,
                    pos['left'],
                    pos['top'],
                    width=pos['width'],
                    height=pos.get('height')
                )
                self._log(f"  Added {chart_name} to slide")
    
    def _add_summary_text(self, slide, statistics: Dict):
        """
        Add summary text box to slide
        
        Args:
            slide: PowerPoint slide object
            statistics: Summary statistics dictionary
        """
        try:
            # Create text box for summary
            left = Inches(0.5)
            top = Inches(0.8)
            width = Inches(9)
            height = Inches(0.6)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Add summary text
            p = text_frame.paragraphs[0]
            p.text = (
                f"Total Risks: {statistics.get('total_risks', 0)} | "
                f"Open: {statistics.get('open_risks', 0)} | "
                f"Closed: {statistics.get('closed_risks', 0)} | "
                f"Business Areas: {statistics.get('business_areas', 0)}"
            )
            p.font.size = Pt(12)
            p.font.bold = True
            
            self._log("  Added summary statistics to slide")
            
        except Exception as e:
            self._log(f"  Error adding summary text: {str(e)}", "warning")
    
    def save_presentation(self) -> bool:
        """
        Save updated presentation
        
        Returns:
            True if successful, False otherwise
        """
        try:
            if self.presentation is None:
                self._log("No presentation to save!", "error")
                return False
            
            # Create output directory if it doesn't exist
            self.output_path.parent.mkdir(parents=True, exist_ok=True)
            
            self._log(f"Saving presentation to: {self.output_path}")
            self.presentation.save(str(self.output_path))
            self._log("Presentation saved successfully")
            return True
            
        except Exception as e:
            self._log(f"Error saving presentation: {str(e)}", "error")
            return False
    
    def create_new_presentation_with_charts(
        self,
        data: pd.DataFrame,
        charts: Dict[str, str],
        statistics: Dict
    ) -> bool:
        """
        Create a new presentation with risk data (if no template exists)
        
        Args:
            data: Risk data DataFrame
            charts: Dictionary of chart names to file paths
            statistics: Summary statistics
            
        Returns:
            True if successful, False otherwise
        """
        try:
            self._log("Creating new presentation with risk data...")
            
            # Create new presentation
            self.presentation = Presentation()
            
            # Add title slide
            title_slide_layout = self.presentation.slide_layouts[0]
            slide = self.presentation.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Technology Performance Review"
            subtitle.text = "Risk Management Report"
            
            # Add risk data slide
            blank_slide_layout = self.presentation.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(blank_slide_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
            title_frame = title_box.text_frame
            title_p = title_frame.paragraphs[0]
            title_p.text = "Risk Management Overview"
            title_p.font.size = Pt(28)
            title_p.font.bold = True
            
            # Add charts and summary
            self._add_charts_to_slide(slide, charts)
            self._add_summary_text(slide, statistics)
            
            # Save
            return self.save_presentation()
            
        except Exception as e:
            self._log(f"Error creating new presentation: {str(e)}", "error")
            return False
